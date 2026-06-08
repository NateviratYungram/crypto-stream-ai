from pathlib import Path

from PIL import Image, ImageDraw, ImageFont


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "CryptoStream_AI_Roadmap_Timeline.png"
PPTX_OUT = ROOT / "CryptoStream_AI_Roadmap_Timeline_Slide.pptx"

W, H = 1920, 1080
BG = (8, 18, 38)
GRID = (11, 27, 51)
CARD = (18, 42, 72)
CARD2 = (22, 52, 86)
WHITE = (245, 248, 252)
MUTED = (170, 185, 205)
CYAN = (38, 198, 218)
BLUE = (66, 133, 244)
PURPLE = (126, 87, 194)
AMBER = (255, 193, 7)
GREEN = (76, 175, 80)
LINE = (84, 126, 165)


def font(size, bold=False):
    candidates = [
        r"C:\Windows\Fonts\LeelawUI.ttf",
        r"C:\Windows\Fonts\arial.ttf",
    ]
    bold_candidates = [
        r"C:\Windows\Fonts\LeelaUIb.ttf",
        r"C:\Windows\Fonts\arialbd.ttf",
    ]
    for path in (bold_candidates if bold else candidates):
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


F_TITLE = font(58, True)
F_SUB = font(27)
F_PHASE = font(27, True)
F_HEAD = font(23, True)
F_BODY = font(21)
F_SMALL = font(17)


def rounded(draw, xy, radius, fill, outline=None, width=2):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def text_center(draw, xy, text, fnt, fill=WHITE, spacing=5):
    x1, y1, x2, y2 = xy
    lines = text.split("\n")
    dims = [draw.textbbox((0, 0), line, font=fnt) for line in lines]
    widths = [b[2] - b[0] for b in dims]
    heights = [b[3] - b[1] for b in dims]
    total_h = sum(heights) + spacing * (len(lines) - 1)
    y = y1 + (y2 - y1 - total_h) / 2
    for line, w, h in zip(lines, widths, heights):
        draw.text((x1 + (x2 - x1 - w) / 2, y), line, font=fnt, fill=fill)
        y += h + spacing


def pill(draw, xy, text, color):
    rounded(draw, xy, 26, color, None, 0)
    text_center(draw, xy, text, F_HEAD, BG)


def arrow(draw, start, end, color=CYAN, width=6):
    draw.line([start, end], fill=color, width=width)
    x2, y2 = end
    draw.polygon([(x2, y2), (x2 - 24, y2 - 14), (x2 - 24, y2 + 14)], fill=color)


img = Image.new("RGB", (W, H), BG)
d = ImageDraw.Draw(img)

for x in range(0, W, 80):
    d.line([(x, 0), (x, H)], fill=GRID, width=1)
for y in range(0, H, 80):
    d.line([(0, y), (W, y)], fill=GRID, width=1)

d.text((95, 72), "CryptoStream AI Roadmap", font=F_TITLE, fill=WHITE)
d.text((98, 143), "Prototype with production-style architecture", font=F_SUB, fill=CYAN)
d.text((98, 184), "From data foundation to stream intelligence, risk-aware workflow and production hardening", font=F_SMALL, fill=MUTED)

timeline_y = 420
d.line([(180, timeline_y), (1740, timeline_y)], fill=LINE, width=6)

phases = [
    {
        "label": "Phase 1",
        "title": "Data Ingestion &\nStorage Foundation",
        "items": ["Source connectors", "Database foundation", "Parquet data lake"],
        "x": 260,
        "color": CYAN,
    },
    {
        "label": "Phase 2",
        "title": "Stream Processing &\nIntelligence Layer",
        "items": ["Kafka & Flink pipeline", "ML signal scoring", "Intelligence API"],
        "x": 710,
        "color": BLUE,
    },
    {
        "label": "Phase 3",
        "title": "Risk Controls,\nRAG & Alerts",
        "items": ["Risk guardrails", "RAG retrieval", "Telegram alerts"],
        "x": 1160,
        "color": PURPLE,
    },
    {
        "label": "Phase 4",
        "title": "Dashboard & Monitoring\nProduction Hardening",
        "items": ["React dashboard", "Prometheus & Grafana", "Readiness checks"],
        "x": 1585,
        "color": GREEN,
    },
]

for i, phase in enumerate(phases):
    x = phase["x"]
    color = phase["color"]
    # Timeline node
    d.ellipse((x - 30, timeline_y - 30, x + 30, timeline_y + 30), fill=BG, outline=color, width=7)
    d.ellipse((x - 12, timeline_y - 12, x + 12, timeline_y + 12), fill=color)

    # Phase pill
    pill(d, (x - 95, 290, x + 95, 348), phase["label"], color)

    # Connector down/up
    d.line([(x, 350), (x, timeline_y - 32)], fill=color, width=4)
    d.line([(x, timeline_y + 32), (x, 520)], fill=color, width=4)

    # Card
    card_w = 360
    card_h = 285
    card_x1 = x - card_w // 2
    card_y1 = 520
    card_x2 = x + card_w // 2
    card_y2 = card_y1 + card_h
    rounded(d, (card_x1, card_y1, card_x2, card_y2), 28, CARD, (54, 91, 128), 2)
    rounded(d, (card_x1, card_y1, card_x1 + 16, card_y2), 8, color)

    text_center(d, (card_x1 + 32, card_y1 + 28, card_x2 - 24, card_y1 + 104), phase["title"], F_PHASE, WHITE, spacing=8)

    item_y = card_y1 + 130
    for item in phase["items"]:
        d.ellipse((card_x1 + 46, item_y + 8, card_x1 + 58, item_y + 20), fill=color)
        d.text((card_x1 + 76, item_y), item, font=F_BODY, fill=MUTED)
        item_y += 42

    if i < len(phases) - 1:
        arrow(d, (x + 60, timeline_y), (phases[i + 1]["x"] - 62, timeline_y), color=CYAN, width=4)

# Bottom legend
rounded(d, (420, 895, 1500, 978), 28, CARD2, (54, 91, 128), 2)
text_center(
    d,
    (450, 910, 1470, 962),
    "Roadmap focus: stable data foundation  >  stream intelligence  >  risk-aware alerts  >  operational hardening",
    F_BODY,
    WHITE,
)

img.save(OUT)
print(OUT)

# Create a single-slide PowerPoint that contains the rendered roadmap image.
try:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(OUT), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.core_properties.title = "CryptoStream AI Roadmap Timeline"
    prs.save(PPTX_OUT)
    print(PPTX_OUT)
except Exception as exc:
    print(f"pptx skipped: {exc}")
