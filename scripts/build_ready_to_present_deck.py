from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "CryptoStream_AI_Project_Pitching_Technical_Deck_v2_Expanded.pptx"
FLOWCHART = ROOT / "CryptoStream_AI_Architecture_Flowchart.png"
ROADMAP = ROOT / "CryptoStream_AI_Roadmap_Timeline.png"
OUT = ROOT / "CryptoStream_AI_Ready_To_Present_Final.pptx"


def add_image_slide(prs: Presentation, image_path: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    return slide


def move_slide(prs: Presentation, old_index: int, new_index: int):
    """Move a slide by manipulating the slide id list. Indexes are zero-based."""
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public reorder API.
    slides = list(sld_id_lst)
    slide_id = slides[old_index]
    sld_id_lst.remove(slide_id)
    sld_id_lst.insert(new_index, slide_id)


def main():
    for path in [SOURCE, FLOWCHART, ROADMAP]:
        if not path.exists():
            raise FileNotFoundError(path)

    prs = Presentation(str(SOURCE))

    add_image_slide(prs, FLOWCHART)
    flowchart_index = len(prs.slides) - 1
    move_slide(prs, flowchart_index, 2)

    add_image_slide(prs, ROADMAP)
    roadmap_index = len(prs.slides) - 1
    move_slide(prs, roadmap_index, 14)

    prs.core_properties.title = "CryptoStream AI - Ready to Present Final Deck"
    prs.core_properties.subject = "AI Market Intelligence and Real-time Data Pipeline"
    prs.core_properties.author = "CryptoStream AI / Software, Data & AI Engineering"
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
