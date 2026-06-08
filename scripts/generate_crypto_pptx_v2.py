from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parents[1] / "CryptoStream_AI_Project_Pitching_Technical_Deck_v2_Expanded.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(8, 18, 38)
CYAN = RGBColor(38, 198, 218)
GREEN = RGBColor(76, 175, 80)
AMBER = RGBColor(255, 193, 7)
RED = RGBColor(239, 83, 80)
WHITE = RGBColor(245, 248, 252)
MUTED = RGBColor(168, 183, 202)
CARD = RGBColor(19, 42, 72)
CARD2 = RGBColor(24, 52, 88)
PURPLE = RGBColor(126, 87, 194)
BLUE = RGBColor(66, 133, 244)
TEXT = RGBColor(215, 225, 238)
FONT = "Leelawadee UI"
FONT_LATIN = "Aptos"


def bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY


def add_text(slide, value, x, y, w, h, size=14, color=WHITE, bold=False, align="left", font=FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    p = frame.paragraphs[0]
    p.text = value
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    return box


def title(slide, value, subtitle, num):
    add_text(slide, value, 0.55, 0.28, 8.7, 0.48, 22, WHITE, True)
    add_text(slide, subtitle, 0.57, 0.78, 10.8, 0.34, 9.5, MUTED)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.18), Inches(2.4), Inches(0.035))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()
    add_text(slide, f"CryptoStream AI | {num:02d}", 11.3, 0.23, 1.5, 0.25, 8, MUTED, False, "right")


def note(slide, value):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(6.78), Inches(12.25), Inches(0.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(11, 25, 46)
    shape.line.color.rgb = RGBColor(42, 73, 105)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = "Speaker Notes: " + value
    p.font.name = FONT
    p.font.size = Pt(7.4)
    p.font.color.rgb = RGBColor(190, 204, 222)


def bullets(items):
    return "\n".join(f"• {item}" for item in items)


def card(slide, x, y, w, h, heading, body="", accent=CYAN, heading_size=13, body_size=9.3):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = RGBColor(43, 72, 105)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, heading, x + 0.14, y + 0.08, w - 0.22, 0.31, heading_size, WHITE, True)
    if body:
        add_text(slide, body, x + 0.14, y + 0.43, w - 0.22, h - 0.5, body_size, TEXT)


def pill(slide, x, y, w, h, value, accent=CYAN, size=9.5, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD2
    shape.line.color.rgb = accent
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.text = value
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = WHITE


def arrow(slide, x1, y1, x2, y2, color=CYAN, width=1.8):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)


def new_slide(num, heading, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, heading, subtitle, num)
    return s


# 1 Title
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
add_text(s, "CryptoStream AI", 0.7, 1.05, 5.9, 0.75, 38, WHITE, True)
add_text(s, "From Noisy Market Data to Actionable Intelligence", 0.73, 1.92, 6.3, 0.35, 18, CYAN, True, font=FONT_LATIN)
add_text(s, "AI Market Intelligence & Real-time Data Platform", 0.75, 2.47, 5.8, 0.35, 15, TEXT, True)
add_text(s, "ยกระดับการตัดสินใจด้านการลงทุนด้วยข้อมูลเรียลไทม์, AI Context และ Risk Guardrails ในแพลตฟอร์มเดียว", 0.75, 2.98, 5.95, 0.9, 17)
card(s, 0.75, 4.35, 4.45, 0.9, "Software, Data & AI Engineering Team", "Prototype with Production-style Architecture", GREEN, 12, 10)
for label, x, y, col in [
    ("Streaming\nMarket Data", 7.25, 1.25, CYAN),
    ("AI\nIntelligence", 9.35, 2.75, PURPLE),
    ("Risk\nGuardrails", 7.45, 4.45, AMBER),
    ("Dashboard\nAlerts", 10.85, 4.2, GREEN),
]:
    oval = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(1.55), Inches(1.55))
    oval.fill.solid()
    oval.fill.fore_color.rgb = RGBColor(16, 41, 70)
    oval.line.color.rgb = col
    oval.line.width = Pt(2)
    frame = oval.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_LATIN
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
arrow(s, 8.55, 2.0, 9.35, 3.1)
arrow(s, 10.55, 3.35, 11.1, 4.2, PURPLE)
arrow(s, 8.15, 4.45, 9.55, 3.8, AMBER)
add_text(s, "CryptoStream AI | 01", 11.3, 0.23, 1.5, 0.25, 8, MUTED, False, "right")
note(s, "เปิดด้วยภาพรวมว่า CryptoStream AI เชื่อม data pipeline, AI intelligence และ risk-aware workflow เข้าด้วยกัน ไม่ใช่เพียง dashboard")

# 2 Overview
s = new_slide(2, "Project Overview", "ระบบทำอะไร รับข้อมูลอะไร และส่งผลลัพธ์ให้ใคร")
for value, x, c in [
    ("Market Data\nBinance / Yahoo / FRED / MT5", 0.85, CYAN),
    ("Data Pipeline\nStreaming + Batch", 3.55, BLUE),
    ("AI Intelligence\nSignal + Context", 6.05, PURPLE),
    ("Risk Guardrails\nFreshness + Confidence", 8.55, AMBER),
    ("Decision Support\nDashboard / Alert / Chat", 10.9, GREEN),
]:
    pill(s, x, 2.1, 2.0, 0.85, value, c, 8.4)
for a, b in [(2.85, 3.55), (5.55, 6.05), (8.05, 8.55), (10.55, 10.9)]:
    arrow(s, a, 2.52, b, 2.52)
card(s, 0.85, 3.75, 3.45, 1.55, "Input", bullets(["Crypto price / volume / candles", "Equities & indices", "Macro indicators", "Broker/API data"]), CYAN)
card(s, 4.9, 3.75, 3.45, 1.55, "Processing", bullets(["Ingestion, transformation", "Signal scoring", "RAG context retrieval", "Risk checks"]), PURPLE)
card(s, 8.95, 3.75, 3.45, 1.55, "Output", bullets(["Insight summary", "Risk level", "Telegram alert", "Operator workflow"]), GREEN)
note(s, "สไลด์นี้ทำหน้าที่ตอบคำถามแรกของผู้ฟัง: ระบบรับข้อมูลตลาดหลายแหล่ง ประมวลผลด้วย pipeline และ AI แล้วส่งเป็น insight, risk และ alert ให้ผู้ใช้")

# 3 Pain
s = new_slide(3, "Pain Points & Solution", "จากข้อมูลกระจัดกระจาย สู่ insight ที่พร้อมใช้ใน workflow เดียว")
card(s, 0.65, 1.55, 3.0, 2.35, "Before", bullets(["ข้อมูลตลาดกระจัดกระจายหลายแหล่ง", "วิเคราะห์บริบทช้าเมื่อเกิด market movement", "ขาด risk visibility ก่อนตัดสินใจ", "ไม่มี workflow กลางสำหรับ analyst/operator"]), RED)
card(s, 9.65, 1.55, 3.0, 2.35, "After", bullets(["รวม streaming และ batch ใน pipeline เดียว", "AI ช่วยสรุป context และ signal", "มี risk guardrails และ anomaly detection", "ส่ง insight ผ่าน dashboard, alert และ workflow"]), GREEN)
for value, x, col in [
    ("Market Movement\nDetected", 1.05, CYAN),
    ("Analyst วิเคราะห์\nไม่ทันต่อเหตุการณ์", 4.0, AMBER),
    ("CryptoStream AI\nsignals + risk context", 6.95, PURPLE),
    ("Faster & safer\ndecision making", 10.15, GREEN),
]:
    pill(s, x, 4.55, 2.35, 0.75, value, col, 10)
for a, b in [(3.4, 4.0), (6.35, 6.95), (9.3, 10.15)]:
    arrow(s, a, 4.92, b, 4.92)
note(s, "ใช้ BTC market movement เป็นเรื่องเล่า: เมื่อเหตุการณ์เกิดเร็ว ระบบต้องช่วยรวมข้อมูลและบริบทให้ทันกับเวลาตัดสินใจ")

# 4 Users
s = new_slide(4, "Target Users & Use Cases", "ผู้ใช้แต่ละกลุ่มเห็นข้อมูลเดียวกัน แต่ได้มุมมองที่เหมาะกับหน้าที่")
for i, (h, b, c) in enumerate([
    ("Trader", "Signal และ alert ที่ตอบสนองเร็ว\nใช้ประกอบจังหวะเข้า/ออกตลาด", CYAN),
    ("Market Analyst", "Market context, historical pattern\nและข้อมูลประกอบการวิเคราะห์", PURPLE),
    ("Operator", "Pipeline health, monitoring\nและ incident alert", AMBER),
    ("Stakeholder", "ภาพรวม risk, system performance\nและ business value", GREEN),
]):
    card(s, 0.7 + i * 3.12, 1.55, 2.75, 1.65, h, b, c, 14, 9.5)
add_text(s, "Core Use Cases", 0.75, 3.62, 2.5, 0.35, 17, WHITE, True)
for i, u in enumerate(["Real-time market monitoring", "AI-assisted context summary", "Anomaly detection & alerting", "Risk-aware decision support", "Pipeline observability"]):
    pill(s, 1.0 + (i % 3) * 4.05, 4.15 + (i // 3) * 0.9, 3.4, 0.48, u, [CYAN, PURPLE, AMBER, GREEN, BLUE][i], 9.3, False)
note(s, "ระบบช่วยให้หลายทีมใช้ source of truth เดียวกัน แต่แสดงข้อมูลตามบทบาทของแต่ละทีม")

# 5 User Workflow
s = new_slide(5, "Actual User Workflow", "เส้นทางใช้งานจริงของ Analyst / Operator")
steps = [
    ("Open Dashboard", "เลือก asset เช่น BTC", CYAN),
    ("View Signal", "เห็น anomaly / price movement", BLUE),
    ("Open AI Context", "อ่านสรุป market context", PURPLE),
    ("Check Risk", "ดู freshness, volatility, confidence", AMBER),
    ("Receive Alert", "Telegram / monitoring channel", GREEN),
    ("Decide", "monitor / escalate / review exposure", GREEN),
]
for i, (h, b, c) in enumerate(steps):
    x = 0.7 + (i % 3) * 4.1
    y = 1.65 + (i // 3) * 1.65
    card(s, x, y, 3.3, 1.15, h, b, c, 13, 9)
    if i in [0, 1, 3, 4]:
        arrow(s, x + 3.3, y + 0.58, x + 3.75, y + 0.58)
arrow(s, 10.7, 2.23, 1.0, 3.88, CYAN, 1.2)
note(s, "สไลด์นี้ทำให้ผู้ฟังเห็นว่าระบบถูกใช้จริงอย่างไร ตั้งแต่เปิด dashboard จนถึงตัดสินใจหรือส่งต่อ workflow")

# 6 Data Sources
s = new_slide(6, "Data Sources & Data Types", "แหล่งข้อมูลแต่ละประเภทถูกใช้เพื่อสร้างบริบทตลาดที่ครบขึ้น")
sources = [
    ("Binance", "Crypto price, volume, candles\nReal-time market movement", CYAN),
    ("Yahoo Finance", "Stocks, indices, historical OHLCV\nCross-market context", BLUE),
    ("FRED", "Macro indicators\nInterest rate, inflation, economic data", PURPLE),
    ("MT5", "Broker/trading terminal integration\nExecution-aware context", AMBER),
    ("Telegram", "Alert delivery channel\nOperator notification", GREEN),
]
for i, (h, b, c) in enumerate(sources):
    card(s, 0.75 + (i % 3) * 4.05, 1.55 + (i // 3) * 1.8, 3.45, 1.25, h, b, c, 14, 8.7)
card(s, 8.85, 3.35, 3.45, 1.25, "Data Quality Focus", bullets(["Freshness", "Completeness", "Source health", "Fallback readiness"]), RED, 14, 8.7)
note(s, "อธิบายให้เห็นว่า data source แต่ละตัวมีหน้าที่ต่างกัน ไม่ใช่ดึงข้อมูลมาโชว์เฉยๆ แต่ใช้สร้าง market context และตรวจ risk")

# 7 Scope
s = new_slide(7, "Project Scope & Key Features", "ขอบเขตชัดเจน แบ่งเป็น 3 โมดูลหลักที่ต่อกันเป็นระบบ end-to-end")
for i, (h, items, c) in enumerate([
    ("Core Data Pipeline", ["Binance, Yahoo Finance, FRED, MT5", "Streaming + batch ingestion", "Kafka, Airflow, pipeline services", "Data freshness + pipeline status"], CYAN),
    ("Intelligence Engine", ["Market signal scoring", "Anomaly detection", "RAG retrieval for context", "Risk guardrails before insight"], PURPLE),
    ("UX & Integrations", ["React dashboard", "AI chat interface", "Grafana monitoring", "Telegram alert + operator workflow"], GREEN),
]):
    card(s, 0.8 + i * 4.15, 1.75, 3.55, 3.15, h, bullets(items), c, 14, 10.3)
for a, b in [(4.35, 4.95), (8.5, 9.1)]:
    arrow(s, a, 3.32, b, 3.32)
add_text(s, "Data Pipeline → Intelligence Engine → User Experience & Integrations", 1.25, 5.45, 10.8, 0.35, 13, TEXT, True, "center")
note(s, "ขอบเขตแบ่งเป็นสามโมดูลเพื่อให้เห็น feasibility และ separation of concerns")

# 8 Architecture
s = new_slide(8, "System Architecture", "Prototype with production-style architecture: scalable, observable, extensible")
for value, x, c in [("Client / React + Vite", 0.75, CYAN), ("API / Chat Server", 3.25, BLUE), ("Intelligence Services", 5.75, PURPLE), ("Storage Layer", 8.55, GREEN), ("Monitoring", 10.95, AMBER)]:
    pill(s, x, 1.45, 2.0, 0.52, value, c, 9.2)
for a, b in [(2.75, 3.25), (5.25, 5.75), (8.0, 8.55), (10.75, 10.95)]:
    arrow(s, a, 1.72, b, 1.72)
card(s, 0.75, 2.55, 2.5, 1.9, "External Services", "• Binance\n• Yahoo Finance\n• FRED\n• MT5\n• Telegram", CYAN, 12, 9.5)
card(s, 3.75, 2.55, 2.8, 1.9, "Streaming & Batch", "• Kafka\n• Flink\n• Airflow\n• dbt\n• Parquet Data Lake", BLUE, 12, 9.5)
card(s, 7.05, 2.55, 2.55, 1.9, "Operational Stores", "• PostgreSQL\n• pgvector\n• Redis\n• SQLite", GREEN, 12, 9.5)
card(s, 10.1, 2.55, 2.45, 1.9, "Observability", "• Prometheus\n• Grafana\n• Alerting\n• Readiness checks", AMBER, 12, 9.5)
for i, (h, b, c) in enumerate([("Scalability", "Kafka/Flink + service separation", CYAN), ("Reliability", "Airflow/dbt + failure alerting", GREEN), ("Security Direction", "API boundary + auth extension", AMBER), ("Separation", "Frontend ไม่ผูกกับ processing", PURPLE)]):
    pill(s, 0.9 + i * 3.05, 5.25, 2.55, 0.52, f"{h}: {b}", c, 8.2, False)
note(s, "Architecture นี้ยังเป็น prototype แต่เดินไปในทิศทาง production มี layer แยกชัดและ harden ต่อได้")

# 9 Data flow
s = new_slide(9, "Data Flow & Pipeline", "ตั้งแต่ source ถึง decision support: end-to-end data platform ไม่ใช่แค่ dashboard")
for i, (h, items, c) in enumerate([
    ("1. Data Ingestion", ["Binance", "Yahoo Finance", "FRED", "MT5", "Airflow", "Kafka"], CYAN),
    ("2. Processing", ["Flink", "Python Intelligence", "ML Scoring", "dbt"], PURPLE),
    ("3. Storage", ["PostgreSQL", "pgvector", "Parquet", "Redis", "SQLite"], GREEN),
    ("4. Decision Support", ["React Dashboard", "Alerts", "Grafana", "AI Chat", "Telegram"], AMBER),
]):
    card(s, 0.65 + i * 3.17, 1.65, 2.65, 3.3, h, bullets(items), c, 13, 10)
    if i < 3:
        arrow(s, 3.3 + i * 3.17, 3.25, 3.82 + i * 3.17, 3.25)
add_text(s, "Flow: Source → Ingest → Process → Store → Model / Analyze → Serve / Alert", 1.0, 5.3, 11.2, 0.35, 15, WHITE, True, "center")
note(s, "Value เกิดจาก pipeline ทั้งเส้น ข้อมูลถูก ingest, process, store, analyze และ serve ไปยังช่องทางตัดสินใจ")

# 10 AI methodology
s = new_slide(10, "AI Intelligence Methodology", "อธิบายว่า AI ในระบบทำงานอย่างไรแบบไม่กล่าวอ้างเกินจริง")
methods = [
    ("Rule-based Checks", "ตรวจ condition พื้นฐาน เช่น missing data, stale data, threshold", CYAN),
    ("Statistical Anomaly", "ตรวจความผิดปกติจากราคา volume และ volatility", BLUE),
    ("ML Signal Scoring", "ให้คะแนน signal confidence เพื่อจัดลำดับความสำคัญ", PURPLE),
    ("RAG Retrieval", "ดึง market notes, macro context และ knowledge base", AMBER),
    ("Human-in-the-loop", "ให้ analyst/operator review ก่อน action สำคัญ", GREEN),
]
for i, (h, b, c) in enumerate(methods):
    card(s, 0.8 + (i % 3) * 4.05, 1.55 + (i // 3) * 1.75, 3.45, 1.25, h, b, c, 13, 8.7)
card(s, 8.9, 3.3, 3.45, 1.35, "AI Output", bullets(["Insight summary", "Confidence score", "Risk explanation", "Recommended next action"]), GREEN, 13, 8.7)
note(s, "สไลด์นี้ลดความกำกวมของคำว่า AI โดยแยก method เป็น rule, statistical, ML scoring, RAG และ human review")

# 11 Risk logic
s = new_slide(11, "Risk Guardrails Logic", "จุดคุมความเสี่ยงก่อนส่ง insight หรือ alert ไปยังผู้ใช้")
checks = [
    ("Data Freshness", "ข้อมูลล่าสุดเกินเวลาที่กำหนดหรือไม่", CYAN),
    ("Missing Data", "มี field สำคัญหายหรือ source ล้มเหลวหรือไม่", BLUE),
    ("Volatility Threshold", "ราคา/volatility สูงผิดปกติหรือไม่", AMBER),
    ("Signal Confidence", "confidence สูงพอสำหรับ alert หรือไม่", PURPLE),
    ("API Health", "external API หรือ broker integration พร้อมใช้งานหรือไม่", RED),
]
for i, (h, b, c) in enumerate(checks):
    card(s, 0.7 + (i % 3) * 4.05, 1.55 + (i // 3) * 1.65, 3.45, 1.15, h, b, c, 13, 8.5)
for value, x, c in [("Low", 1.0, GREEN), ("Medium", 3.0, CYAN), ("High", 5.0, AMBER), ("Critical", 7.0, RED), ("Human Review", 9.0, PURPLE)]:
    pill(s, x, 5.15, 1.55, 0.52, value, c, 9.5)
note(s, "Risk guardrails คือส่วนที่ทำให้ระบบน่าเชื่อถือขึ้น เพราะไม่ส่ง alert จาก signal อย่างเดียว แต่ตรวจ freshness, confidence และ source health ประกอบ")

# 12 Demo
s = new_slide(12, "Demo Scenario: BTC Market Movement Detected", "ตัวอย่าง event จริง: จาก price spike สู่ alert ที่มี context และ risk level")
actors = ["Market", "Pipeline", "AI Engine", "Risk", "Dashboard", "Telegram"]
xs = [0.9, 3.0, 5.1, 7.2, 9.25, 11.25]
for x, a in zip(xs, actors):
    pill(s, x, 1.55, 1.55, 0.45, a, CYAN if a in ["Market", "Pipeline"] else PURPLE if a == "AI Engine" else AMBER if a == "Risk" else GREEN, 8.5)
    line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 0.775), Inches(2.05), Inches(x + 0.775), Inches(5.55))
    line.line.color.rgb = RGBColor(55, 78, 108)
    line.line.width = Pt(0.8)
for start, end, y, label in [(0, 1, 2.35, "BTC volatility spike"), (1, 2, 2.8, "Event + feature payload"), (2, 3, 3.25, "Signal 82% + context"), (3, 4, 3.7, "Risk Level: High"), (4, 5, 4.15, "Alert + summary")]:
    arrow(s, xs[start] + 1.55, y, xs[end], y)
    add_text(s, label, xs[start] + 1.1, y - 0.24, max(1.7, xs[end] - xs[start] - 0.25), 0.22, 7.5, TEXT, False, "center")
card(s, 0.9, 5.65, 11.55, 0.85, "Expected Output Example", "Signal: BTC volatility spike detected | Confidence: 82% | Risk Level: High | Context: increased market volume | Action: review exposure and monitor next candle confirmation", GREEN, 12, 9.2)
note(s, "ใช้สไลด์นี้เป็น demo script แบบสั้น: ระบบตรวจ event วิเคราะห์ signal ตรวจ risk แล้วส่ง alert พร้อม context")

# 13 Roadmap + status
s = new_slide(13, "Roadmap & Implementation Status", "แยกสถานะปัจจุบันออกจากแผนต่อยอด เพื่อให้ feasibility ชัดเจน")
for i, (h, items, c) in enumerate([
    ("Done / Foundation", ["React/Vite frontend direction", "FastAPI-style backend", "SQLite/PostgreSQL storage", "Basic dashboard/API"], GREEN),
    ("Prototype", ["Intelligence engine", "Signal scoring", "RAG retrieval with pgvector", "Telegram alert flow"], CYAN),
    ("In Progress", ["Monitoring readiness", "Data freshness checks", "Pipeline hardening", "Operator workflow"], AMBER),
    ("Future Extension", ["Kafka/Flink scaling", "Auth/access control", "Broker sandbox guardrails", "Production deployment plan"], PURPLE),
]):
    card(s, 0.65 + i * 3.17, 1.55, 2.65, 3.45, h, bullets(items), c, 12.5, 8.5)
add_text(s, "Project status: Prototype with production-style architecture, not full production-ready system.", 1.0, 5.45, 11.2, 0.35, 13, TEXT, True, "center")
note(s, "สไลด์นี้สำคัญสำหรับกรรมการหรือ stakeholder เพราะบอกตรงๆ ว่าอะไรทำแล้ว อะไรยังเป็น prototype และอะไรเป็น roadmap")

# 14 Metrics risks limitations
s = new_slide(14, "Metrics, Risks, Limitations & Q&A", "Success metrics ที่วัดได้ พร้อมข้อจำกัดและแผนลดความเสี่ยง")
for i, (num, label, c) in enumerate([
    ("1-3 นาที", "Initial analysis time\nจากเดิม 15-30 นาที", CYAN),
    ("5-30 วินาที", "Real-time data freshness\nตามประเภท source", GREEN),
    ("5-10 วินาที", "Anomaly detection latency\nหลัง event เข้าระบบ", AMBER),
    ("95-99%", "Prototype pipeline\navailability target", PURPLE),
]):
    card(s, 0.65 + i * 3.15, 1.35, 2.65, 1.0, num, label, c, 18, 8)
card(s, 0.75, 2.72, 3.8, 1.55, "Key Risks", bullets(["Upstream data instability", "False positives", "External API rate limits", "Pipeline failure/stale data"]), RED, 13, 8.5)
card(s, 4.75, 2.72, 3.8, 1.55, "Mitigation", bullets(["Fallback source + retry", "Confidence score + thresholds", "Health checks + alerting", "Human review workflow"]), GREEN, 13, 8.5)
card(s, 8.75, 2.72, 3.8, 1.55, "Limitations", bullets(["ยังเป็น prototype", "ต้อง validate signal เพิ่ม", "ยังต้อง harden security/deployment", "broker action ต้องใช้ sandbox ก่อน"]), AMBER, 13, 8.5)
card(s, 0.75, 5.15, 11.8, 0.85, "Q&A", "CryptoStream AI is designed as a scalable foundation for real-time market intelligence, risk-aware analytics, and operational decision support.\nTeam: CryptoStream AI Development Team | Role: Software, Data & AI Engineering | Contact: [ใส่อีเมลหรือช่องทางติดต่อ]", CYAN, 15, 8.7)
note(s, "ปิดด้วย metrics ที่วัดได้ ความเสี่ยง ข้อจำกัด และย้ำว่าโปรเจกต์ต่อยอดได้ทั้งเชิงธุรกิจและเชิงเทคนิค")

prs.core_properties.title = "CryptoStream AI - Expanded Project Pitching & Technical Deck"
prs.core_properties.subject = "AI Market Intelligence and Real-time Data Pipeline"
prs.core_properties.author = "CryptoStream AI / Software, Data & AI Engineering"
prs.save(OUT)
print(OUT)
