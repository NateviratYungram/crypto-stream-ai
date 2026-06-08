from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parents[1] / "CryptoStream_AI_Project_Pitching_Technical_Deck_UTF8.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(8, 18, 38)
CYAN = RGBColor(38, 198, 218)
GREEN = RGBColor(76, 175, 80)
AMBER = RGBColor(255, 193, 7)
RED = RGBColor(239, 83, 80)
WHITE = RGBColor(245, 248, 252)
MUTED = RGBColor(168, 183, 202)
CARD = RGBColor(19, 42, 72)
CARD2 = RGBColor(24, 52, 88)
PURPLE = RGBColor(126, 87, 194)
BLUE = RGBColor(66, 133, 244)
TEXT = RGBColor(215, 225, 238)
FONT = "Leelawadee UI"
FONT_LATIN = "Aptos"


def bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY


def text(slide, value, x, y, w, h, size=14, color=WHITE, bold=False, align="left", font=FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    p = frame.paragraphs[0]
    p.text = value
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    return box


def title(slide, value, subtitle, num):
    text(slide, value, 0.55, 0.28, 8.7, 0.48, 22, WHITE, True)
    text(slide, subtitle, 0.57, 0.78, 10.8, 0.34, 9.5, MUTED)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.18), Inches(2.4), Inches(0.035))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()
    text(slide, f"CryptoStream AI | {num:02d}", 11.3, 0.23, 1.5, 0.25, 8, MUTED, False, "right")


def note(slide, value):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(6.78), Inches(12.25), Inches(0.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(11, 25, 46)
    shape.line.color.rgb = RGBColor(42, 73, 105)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = "Speaker Notes: " + value
    p.font.name = FONT
    p.font.size = Pt(7.6)
    p.font.color.rgb = RGBColor(190, 204, 222)


def card(slide, x, y, w, h, heading, body="", accent=CYAN, heading_size=13, body_size=9.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = RGBColor(43, 72, 105)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    text(slide, heading, x + 0.14, y + 0.08, w - 0.22, 0.31, heading_size, WHITE, True)
    if body:
        text(slide, body, x + 0.14, y + 0.43, w - 0.22, h - 0.5, body_size, TEXT)


def pill(slide, x, y, w, h, value, accent=CYAN, size=10, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD2
    shape.line.color.rgb = accent
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.text = value
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = WHITE


def arrow(slide, x1, y1, x2, y2, color=CYAN, width=1.8):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)


def bullets(items):
    return "\n".join(f"• {item}" for item in items)


# Slide 1
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
text(s, "CryptoStream AI", 0.7, 1.15, 5.8, 0.75, 38, WHITE, True)
text(s, "From Noisy Market Data to Actionable Intelligence", 0.73, 2.0, 6.2, 0.35, 18, CYAN, True, font=FONT_LATIN)
text(s, "AI Market Intelligence & Real-time Data Platform", 0.75, 2.55, 5.8, 0.35, 15, TEXT, True)
text(s, "ยกระดับการตัดสินใจด้านการลงทุนด้วยข้อมูลเรียลไทม์, AI Context และ Risk Guardrails ในแพลตฟอร์มเดียว", 0.75, 3.05, 5.9, 0.9, 17)
card(s, 0.75, 4.45, 4.3, 0.85, "Software, Data & AI Engineering Team", "Prototype with Production-style Architecture", GREEN, 12, 10)
for label, x, y, col in [
    ("Streaming\nMarket Data", 7.25, 1.25, CYAN),
    ("AI\nIntelligence", 9.35, 2.75, PURPLE),
    ("Risk\nGuardrails", 7.45, 4.45, AMBER),
    ("Dashboard\nAlerts", 10.85, 4.2, GREEN),
]:
    oval = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(1.55), Inches(1.55))
    oval.fill.solid()
    oval.fill.fore_color.rgb = RGBColor(16, 41, 70)
    oval.line.color.rgb = col
    oval.line.width = Pt(2)
    frame = oval.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_LATIN
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
arrow(s, 8.55, 2.0, 9.35, 3.1)
arrow(s, 10.55, 3.35, 11.1, 4.2, PURPLE)
arrow(s, 8.15, 4.45, 9.55, 3.8, AMBER)
text(s, "CryptoStream AI | 01", 11.3, 0.23, 1.5, 0.25, 8, MUTED, False, "right")
note(s, "เปิดด้วยภาพรวมว่า CryptoStream AI ไม่ใช่เพียง dashboard แต่เป็น data platform ที่เชื่อม pipeline, AI intelligence และ risk-aware workflow เข้าด้วยกัน")

# Slide 2
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
title(s, "Pain Points & Solution", "จากข้อมูลกระจัดกระจาย สู่ insight ที่พร้อมใช้ใน workflow เดียว", 2)
card(s, 0.65, 1.55, 3.0, 2.35, "Before", bullets(["ข้อมูลตลาดกระจัดกระจายหลายแหล่ง", "วิเคราะห์บริบทช้าเมื่อเกิด market movement", "ขาด risk visibility ก่อนตัดสินใจ", "ไม่มี workflow กลางสำหรับ analyst/operator"]), RED)
card(s, 9.65, 1.55, 3.0, 2.35, "After", bullets(["รวม streaming และ batch ใน pipeline เดียว", "AI ช่วยสรุป context และ signal", "มี risk guardrails และ anomaly detection", "ส่ง insight ผ่าน dashboard, alert และ workflow"]), GREEN)
for value, x, col in [
    ("Market Movement\nDetected", 1.05, CYAN),
    ("Analyst วิเคราะห์\nไม่ทันต่อเหตุการณ์", 4.0, AMBER),
    ("CryptoStream AI\naggregates signals + risk", 6.95, PURPLE),
    ("Faster & safer\ndecision making", 10.15, GREEN),
]:
    pill(s, x, 4.55, 2.35, 0.75, value, col, 10)
for a, b in [(3.4, 4.0), (6.35, 6.95), (9.3, 10.15)]:
    arrow(s, a, 4.92, b, 4.92)
text(s, "Diagram: User Journey Flowchart", 0.8, 4.1, 4.5, 0.3, 10, MUTED, True)
note(s, "ใช้ demo scenario BTC Market Movement Detected เพื่อเล่าให้เห็น pain point: ข้อมูลมีหลายแหล่งและเวลาตัดสินใจสั้น ระบบจึงรวม signal, context และ risk guardrails ให้พร้อมใช้งาน")

# Slide 3
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
title(s, "Target Users & Use Cases", "ผู้ใช้แต่ละกลุ่มเห็นข้อมูลเดียวกัน แต่ได้มุมมองที่เหมาะกับหน้าที่", 3)
for i, (h, b, c) in enumerate([
    ("Trader", "Signal และ alert ที่ตอบสนองเร็ว\nใช้ประกอบจังหวะเข้า/ออกตลาด", CYAN),
    ("Market Analyst", "Market context, historical pattern\nและข้อมูลประกอบการวิเคราะห์", PURPLE),
    ("Operator", "Pipeline health, monitoring\nและ incident alert", AMBER),
    ("Stakeholder", "ภาพรวม risk, system performance\nและ business value", GREEN),
]):
    card(s, 0.7 + i * 3.12, 1.55, 2.75, 1.65, h, b, c, 14, 9.5)
text(s, "Core Use Cases", 0.75, 3.62, 2.5, 0.35, 17, WHITE, True)
for i, u in enumerate(["Real-time market movement monitoring", "AI-assisted market context summary", "Anomaly detection and alerting", "Risk-aware decision support", "Pipeline observability and operational dashboard"]):
    pill(s, 1.0 + (i % 3) * 4.05, 4.15 + (i // 3) * 0.9, 3.4, 0.48, u, [CYAN, PURPLE, AMBER, GREEN, BLUE][i], 9.3, False)
note(s, "ระบบไม่ได้ออกแบบเพื่อคนกลุ่มเดียว แต่ทำให้ Business, Data และ Engineering ใช้ source of truth ร่วมกันได้ โดยแสดงผลตามบริบทของบทบาท")

# Slide 4
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
title(s, "Project Scope & Key Features", "ขอบเขตชัดเจน แบ่งเป็น 3 โมดูลหลักที่ต่อกันเป็นระบบ end-to-end", 4)
for i, (h, items, c) in enumerate([
    ("Core Data Pipeline", ["Binance, Yahoo Finance, FRED, MT5", "Streaming + batch ingestion", "Kafka, Airflow, pipeline services", "Data freshness + pipeline status"], CYAN),
    ("Intelligence Engine", ["Market signal scoring", "Anomaly detection", "RAG retrieval for context", "Risk guardrails before insight"], PURPLE),
    ("UX & Integrations", ["React dashboard", "AI chat interface", "Grafana monitoring", "Telegram alert + operator workflow"], GREEN),
]):
    card(s, 0.8 + i * 4.15, 1.75, 3.55, 3.15, h, bullets(items), c, 14, 10.3)
for a, b in [(4.35, 4.95), (8.5, 9.1)]:
    arrow(s, a, 3.32, b, 3.32)
text(s, "Functional Block Diagram", 0.9, 5.25, 3.2, 0.28, 10, MUTED, True)
text(s, "Data Pipeline ส่งข้อมูลที่พร้อมใช้ให้ Intelligence Engine จากนั้นส่ง insight, alert และ workflow ไปยังผู้ใช้และระบบปฏิบัติการ", 0.9, 5.58, 11.2, 0.42, 12, TEXT)
note(s, "โปรเจกต์ทำได้จริงเพราะแบ่งเป็นโมดูลชัดเจน: pipeline สร้างข้อมูลที่เชื่อถือได้, intelligence layer สร้าง insight, และ UX/integration ทำให้ insight ถูกใช้งานได้")

# Slide 5
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
title(s, "System Architecture", "Prototype with production-style architecture: scalable, observable, extensible", 5)
for value, x, c in [("Client / React + Vite", 0.75, CYAN), ("API / Chat Server", 3.25, BLUE), ("Intelligence Services", 5.75, PURPLE), ("Storage Layer", 8.55, GREEN), ("Monitoring", 10.95, AMBER)]:
    pill(s, x, 1.45, 2.0, 0.52, value, c, 9.2)
for a, b in [(2.75, 3.25), (5.25, 5.75), (8.0, 8.55), (10.75, 10.95)]:
    arrow(s, a, 1.72, b, 1.72)
card(s, 0.75, 2.55, 2.5, 1.9, "External Services", "• Binance\n• Yahoo Finance\n• FRED\n• MT5\n• Telegram", CYAN, 12, 9.5)
card(s, 3.75, 2.55, 2.8, 1.9, "Streaming & Batch", "• Kafka\n• Flink\n• Airflow\n• dbt\n• Parquet Data Lake", BLUE, 12, 9.5)
card(s, 7.05, 2.55, 2.55, 1.9, "Operational Stores", "• PostgreSQL\n• pgvector\n• Redis\n• SQLite", GREEN, 12, 9.5)
card(s, 10.1, 2.55, 2.45, 1.9, "Observability", "• Prometheus\n• Grafana\n• Alerting\n• Readiness checks", AMBER, 12, 9.5)
for i, (h, b, c) in enumerate([("Scalability", "Kafka/Flink + service separation", CYAN), ("Reliability", "Airflow/dbt + failure alerting", GREEN), ("Security Direction", "API boundary + secret/auth extension", AMBER), ("Separation", "Frontend ไม่ผูกกับ processing", PURPLE)]):
    pill(s, 0.9 + i * 3.05, 5.25, 2.55, 0.52, f"{h}: {b}", c, 8.2, False)
note(s, "Architecture นี้ยังเป็น prototype แต่เดินไปในทิศทาง production: มี layer แยกชัด รองรับ streaming/batch/AI/storage/monitoring และสามารถ harden ต่อได้")

# Slide 6
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
title(s, "Data Flow & Pipeline", "ตั้งแต่ source ถึง decision support: end-to-end data platform ไม่ใช่แค่ dashboard", 6)
for i, (h, items, c) in enumerate([
    ("1. Data Ingestion", ["Binance", "Yahoo Finance", "FRED", "MT5", "Airflow", "Kafka"], CYAN),
    ("2. Processing", ["Flink", "Python Intelligence", "ML Scoring", "dbt"], PURPLE),
    ("3. Storage", ["PostgreSQL", "pgvector", "Parquet", "Redis", "SQLite"], GREEN),
    ("4. Decision Support", ["React Dashboard", "Alerts", "Grafana", "AI Chat", "Telegram"], AMBER),
]):
    card(s, 0.65 + i * 3.17, 1.65, 2.65, 3.3, h, bullets(items), c, 13, 10)
    if i < 3:
        arrow(s, 3.3 + i * 3.17, 3.25, 3.82 + i * 3.17, 3.25)
text(s, "Flow: Source → Ingest → Process → Store → Model / Analyze → Serve / Alert", 1.0, 5.3, 11.2, 0.35, 15, WHITE, True, "center")
note(s, "Value เกิดจาก pipeline ทั้งเส้น ไม่ใช่หน้า dashboard เพียงอย่างเดียว ข้อมูลถูก ingest, process, store, analyze และ serve ออกไปยังช่องทางที่ใช้ตัดสินใจ")

# Slide 7
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
title(s, "AI Intelligence & Risk Guardrails", "AI ช่วยจัดลำดับความสำคัญ แต่มี guardrails เพื่อคุมความเสี่ยงและความน่าเชื่อถือ", 7)
for value, x, w, c in [("Raw Signals", 0.85, 2.3, CYAN), ("AI Scoring\nSignal + Anomaly", 3.25, 2.3, PURPLE), ("Risk Guardrails\nFreshness + Confidence", 5.65, 2.8, AMBER), ("Actionable Insight\nRisk + Evidence + Action", 8.65, 3.45, GREEN)]:
    pill(s, x, 2.15, w, 0.82, value, c, 10)
for a, b in [(3.15, 3.25), (5.55, 5.65), (8.45, 8.65)]:
    arrow(s, a, 2.56, b, 2.56)
card(s, 0.9, 3.65, 5.45, 1.65, "AI Intelligence Layer", bullets(["Market signal scoring", "Anomaly detection", "Context summarization", "RAG-based knowledge retrieval", "Confidence scoring"]), PURPLE, 13, 9.5)
card(s, 6.85, 3.65, 5.45, 1.65, "Risk Guardrails", bullets(["ตรวจ data freshness และ missing/stale data", "ตรวจ volatility threshold และ signal confidence", "ลด false positive ด้วย rule-based checks + human review"]), AMBER, 13, 9.5)
text(s, "Output: Insight summary • Risk level • Alert priority • Supporting evidence • Recommended next action", 1.0, 5.75, 11.2, 0.35, 12, TEXT, True, "center")
note(s, "AI ไม่ได้ตัดสินใจแทนมนุษย์ แต่ช่วยสรุป context, ให้ confidence score, ตรวจ risk guardrails และส่งหลักฐานประกอบเพื่อให้ workflow ปลอดภัยขึ้น")

# Slide 8
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
title(s, "Demo Scenario: BTC Market Movement Detected", "ตัวอย่าง event จริง: จาก price spike สู่ alert ที่มี context และ risk level", 8)
actors = ["Market", "Pipeline", "AI Engine", "Risk", "Dashboard", "Telegram"]
xs = [0.9, 3.0, 5.1, 7.2, 9.25, 11.25]
for x, a in zip(xs, actors):
    pill(s, x, 1.55, 1.55, 0.45, a, CYAN if a in ["Market", "Pipeline"] else PURPLE if a == "AI Engine" else AMBER if a == "Risk" else GREEN, 8.5)
    line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 0.775), Inches(2.05), Inches(x + 0.775), Inches(5.55))
    line.line.color.rgb = RGBColor(55, 78, 108)
    line.line.width = Pt(0.8)
for start, end, y, label in [(0, 1, 2.35, "BTC volatility spike"), (1, 2, 2.8, "Event + feature payload"), (2, 3, 3.25, "Signal 82% + context"), (3, 4, 3.7, "Risk Level: High"), (4, 5, 4.15, "Alert + summary")]:
    arrow(s, xs[start] + 1.55, y, xs[end], y)
    text(s, label, xs[start] + 1.1, y - 0.24, max(1.7, xs[end] - xs[start] - 0.25), 0.22, 7.5, TEXT, False, "center")
card(s, 0.9, 5.65, 11.55, 0.85, "Expected Output Example", "Signal: BTC volatility spike detected | Confidence: 82% | Risk Level: High | Context: increased market volume | Action: review exposure and monitor next candle confirmation", GREEN, 12, 9.2)
note(s, "เมื่อ BTC ขยับผิดปกติ ระบบ ingest ข้อมูล วิเคราะห์ signal ดึง context ตรวจ guardrails แล้วแสดงบน dashboard พร้อมส่ง Telegram alert")

# Slide 9
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
title(s, "Roadmap & Timeline", "4 เฟสจาก foundation สู่ production hardening direction", 9)
phases = [
    ("Phase 1", "Data Ingestion &\nStorage Foundation", ["Source connectors", "PostgreSQL schema", "Parquet layout", "Basic API", "Dashboard skeleton"], CYAN),
    ("Phase 2", "Stream Processing &\nIntelligence Layer", ["Kafka/Flink pipeline", "Signal processor", "ML scoring prototype", "Freshness checks", "Internal API"], PURPLE),
    ("Phase 3", "Risk Controls,\nRAG & Alerts", ["Risk guardrails", "pgvector RAG", "Telegram integration", "Alert priority rules", "Human review"], AMBER),
    ("Phase 4", "Dashboard, Monitoring &\nProduction Hardening", ["React dashboard", "Grafana/Prometheus", "Failure alerting", "Readiness checks", "Scaling plan"], GREEN),
]
for i, (phase, h, items, c) in enumerate(phases):
    x = 0.75 + i * 3.05
    pill(s, x, 1.55, 2.35, 0.42, phase, c, 10)
    card(s, x, 2.15, 2.35, 3.15, h, bullets(items), c, 11, 8.6)
    if i < 3:
        arrow(s, x + 2.35, 3.7, x + 3.05, 3.7)
note(s, "Roadmap เริ่มจาก ingestion/storage ก่อน แล้วค่อยเพิ่ม stream intelligence, risk/RAG/alerts และปิดท้ายด้วย monitoring กับ production hardening")

# Slide 10
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
title(s, "Metrics, Risks & Q&A", "Success metrics ที่วัดได้ พร้อม risk mitigation สำหรับ prototype platform", 10)
for i, (num, label, c) in enumerate([
    ("1-3 นาที", "Initial analysis time\nจากเดิม 15-30 นาที", CYAN),
    ("5-30 วินาที", "Real-time data freshness\nตามประเภท source", GREEN),
    ("5-10 วินาที", "Anomaly detection latency\nหลัง event เข้าระบบ", AMBER),
    ("95-99%", "Prototype pipeline\navailability target", PURPLE),
]):
    card(s, 0.65 + i * 3.15, 1.45, 2.65, 1.1, num, label, c, 20, 8.2)
text(s, "Technical Risks & Mitigation", 0.75, 2.9, 4.2, 0.35, 16, WHITE, True)
for i, (risk, mitigation) in enumerate([
    ("Upstream data instability", "Fallback source, retry, freshness monitoring, source failure alert"),
    ("Model / signal false positives", "Confidence score, rule-based guardrails, threshold tuning, human review"),
    ("External API / broker integration", "Adapter pattern, sandbox testing, rate limit handling, health checks"),
    ("Pipeline failure / stale data", "Prometheus/Grafana, Airflow alerting, readiness checks, retry/dead-letter queue"),
]):
    y = 3.35 + i * 0.55
    card(s, 0.75, y, 3.1, 0.44, risk, "", RED if i == 0 else AMBER, 8.6, 8)
    text(s, mitigation, 4.05, y + 0.04, 8.25, 0.33, 8.7, TEXT)
card(s, 0.75, 5.75, 11.8, 0.8, "Q&A", "CryptoStream AI is designed as a scalable foundation for real-time market intelligence, risk-aware analytics, and operational decision support.\nTeam: CryptoStream AI Development Team | Role: Software, Data & AI Engineering | Contact: [ใส่อีเมลหรือช่องทางติดต่อ]", CYAN, 16, 8.8)
note(s, "สรุปว่าระบบมีทั้ง business value และ technical feasibility โดยยังสื่ออย่างตรงไปตรงมาว่าเป็น prototype with production-style architecture ที่ต่อยอดได้")

prs.core_properties.title = "CryptoStream AI - Project Pitching & Technical Deck"
prs.core_properties.subject = "AI Market Intelligence and Real-time Data Pipeline"
prs.core_properties.author = "CryptoStream AI / Software, Data & AI Engineering"
prs.save(OUT)
print(OUT)
